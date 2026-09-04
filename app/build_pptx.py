import io
import os
from functools import lru_cache

from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

MARGIN = Inches(0.6)
TITLE_HEIGHT = Inches(1.0)
GAP = Inches(0.35)

FONT_NAME = "Calibri"
FONTS_DIR = os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
FONT_REGULAR_PATH = os.path.join(FONTS_DIR, "calibri.ttf")
FONT_BOLD_PATH = os.path.join(FONTS_DIR, "calibrib.ttf")

TITLE_COLOR = RGBColor(0x22, 0x22, 0x22)
ACCENT_COLOR = RGBColor(0xC0, 0x39, 0x2B)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)

BOX_INSET_WIDTH_IN = 0.2
BOX_INSET_HEIGHT_IN = 0.1
LINE_SPACING = 1.25
SPACE_AFTER_FACTOR = 0.5
SAFETY_MARGIN = 0.92
BULLET_PREFIX = "•  "


def build_presentation(slides_content, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[6]
    content_top = MARGIN + TITLE_HEIGHT + GAP
    content_height = prs.slide_height - content_top - MARGIN

    for content in slides_content:
        slide = prs.slides.add_slide(layout)
        _add_title(slide, prs.slide_width, content["title"] or "Слайд")
        _add_accent_line(slide, prs.slide_width)

        bullets = content["bullets"] or [content.get("text", "")[:500]]
        images = content.get("images", [])

        if images:
            text_width = int((prs.slide_width - 3 * MARGIN) * 0.55)
            _add_body(slide, MARGIN, content_top, text_width, content_height, bullets)

            picture_left = MARGIN + text_width + MARGIN
            picture_area_width = prs.slide_width - MARGIN - picture_left
            _add_picture_centered(slide, images[0], picture_left, content_top, picture_area_width, content_height)
        else:
            text_width = prs.slide_width - 2 * MARGIN
            _add_body(slide, MARGIN, content_top, text_width, content_height, bullets)

    prs.save(output_path)


def _add_run(paragraph, text, size_pt, color, bold=False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = FONT_NAME
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def _add_title(slide, slide_width, title_text):
    box = slide.shapes.add_textbox(MARGIN, MARGIN, slide_width - 2 * MARGIN, TITLE_HEIGHT)
    box.text_frame.word_wrap = False

    font_size = _pick_title_font_size(title_text, slide_width - 2 * MARGIN)
    _add_run(box.text_frame.paragraphs[0], title_text, font_size, TITLE_COLOR, bold=True)


def _add_accent_line(slide, slide_width):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, MARGIN + TITLE_HEIGHT, slide_width - 2 * MARGIN, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    line.shadow.inherit = False


def _add_body(slide, left, top, width, height, bullets):
    text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
    text_frame.word_wrap = True

    font_size = _pick_body_font_size(bullets, width, height)
    for i, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.line_spacing = LINE_SPACING
        p.space_after = Pt(font_size * SPACE_AFTER_FACTOR)
        _add_run(p, BULLET_PREFIX, font_size, ACCENT_COLOR, bold=True)
        _add_run(p, bullet, font_size, TEXT_COLOR)


@lru_cache(maxsize=None)
def _load_font(path, size_pt):
    try:
        return ImageFont.truetype(path, round(size_pt * 96 / 72))
    except OSError:
        return ImageFont.load_default()


def _wrapped_line_count(text, font, max_width_px):
    words = text.split(" ")
    lines = 1
    current = ""
    for word in words:
        trial = f"{current} {word}".strip()
        if font.getlength(trial) <= max_width_px:
            current = trial
        else:
            lines += 1
            current = word
    return lines


def _usable_width_px(width_emu):
    return max((width_emu / 914400) - BOX_INSET_WIDTH_IN, 0.5) * 96


def _pick_body_font_size(bullets, width_emu, height_emu):
    width_px = _usable_width_px(width_emu)
    height_in = max((height_emu / 914400) - BOX_INSET_HEIGHT_IN, 0.5) * SAFETY_MARGIN

    for size in range(24, 11, -1):
        font = _load_font(FONT_REGULAR_PATH, size)
        ascent, descent = font.getmetrics()
        line_height_in = (ascent + descent) / 96 * LINE_SPACING
        space_after_in = size * SPACE_AFTER_FACTOR / 72

        total_height_in = sum(
            _wrapped_line_count(BULLET_PREFIX + b, font, width_px) * line_height_in + space_after_in
            for b in bullets
        )
        if total_height_in <= height_in:
            return size
    return 12


def _pick_title_font_size(title, width_emu):
    width_px = _usable_width_px(width_emu)
    for size in range(36, 19, -2):
        font = _load_font(FONT_BOLD_PATH, size)
        if font.getlength(title) <= width_px:
            return size
    return 20


def _add_picture_centered(slide, image_bytes, area_left, area_top, area_width, area_height):
    img_width_px, img_height_px = Image.open(io.BytesIO(image_bytes)).size
    aspect = img_width_px / img_height_px

    width = area_width
    height = int(width / aspect)
    if height > area_height:
        height = area_height
        width = int(height * aspect)

    left = area_left + (area_width - width) // 2
    top = area_top + (area_height - height) // 2

    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)
